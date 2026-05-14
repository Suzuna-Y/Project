# app.py (全文)
from flask import Flask, render_template, request, redirect, url_for, session, jsonify
import google.generativeai as genai
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image
import fitz  # PyMuPDF
import io
import re
import base64
import os
from pathlib import Path

app = Flask(__name__)
# セッションキーは環境変数で設定することを推奨（ここは開発用デフォルト）
app.secret_key = os.environ.get("FLASK_SECRET_KEY", "dev-secret-key")

APP_STATE = {
    "question_list": []
}

# ---------------- DEBUG: 環境変数キー検出（マスク表示） ----------------
# もし .env を使うなら読み込む（python-dotenv がある場合）
env_path = Path('.') / '.env'
if env_path.exists():
    try:
        from dotenv import load_dotenv  # type: ignore
        load_dotenv(dotenv_path=env_path)
        print("読み込み: .env ファイルを読み込みました。")
    except Exception:
        print("注意: python-dotenv がインストールされていません（.env を使うなら pip install python-dotenv）")

key_names = ["GENAI_KEY", "GOOGLE_API_KEY", "GOOGLE_APIKEY", "API_KEY"]
found = {}
for name in key_names:
    v = os.environ.get(name)
    if v:
        masked = v if len(v) <= 8 else (v[:4] + "..." + v[-4:])
        print(f"環境変数 {name} = {masked} (length={len(v)})")
        found[name] = v

if not found:
    print("警告: GENAI_KEY / GOOGLE_API_KEY 等の環境変数が見つかりませんでした。")
else:
    API_KEY = found.get("GENAI_KEY") or found.get("GOOGLE_API_KEY") or next(iter(found.values()))
    print("使用する API キーを決定しました（マスク表示）:", (API_KEY[:4] + "..." + API_KEY[-4:]) if len(API_KEY) > 8 else API_KEY)
    # 使いやすいように内部でもセットしておく（安全上の理由で実際の値は表示していません）
    os.environ["__DETECTED_GENAI_KEY__"] = API_KEY

# ---------------- GenAI: 遅延初期化 ----------------


model = None

def ensure_model():
    global model

    API_KEY = (
        os.environ.get("GENAI_KEY")
        or os.environ.get("GOOGLE_API_KEY")
        or os.environ.get("__DETECTED_GENAI_KEY__")
        or ""
    )

    if not API_KEY:
        print("ensure_model: APIキーが見つかりません")
        model = None
        return False

    try:
        genai.configure(api_key=API_KEY)
        # 毎回 fresh な model を作る
        model = genai.GenerativeModel(
            model_name="models/gemini-2.5-flash"
        )
        print("ensure_model: model を新規作成しました")
        return True

    except Exception as e:
        print("ensure_model: 初期化失敗:", repr(e))
        model = None
        return False

# 簡易テスト（開発中のみ：起動時に一度だけ試す）

# ---------------- ユーティリティ: スライド抽出（PPTX） ----------------
def extract_slides_from_pptx_bytes(file_bytes):
    prs = Presentation(io.BytesIO(file_bytes))
    slides = []
    for i, slide in enumerate(prs.slides, start=1):
        slide_text = ""
        images = []
        title = ""
        try:
            if slide.shapes.title and slide.shapes.title.text:
                title = slide.shapes.title.text.strip()
        except Exception:
            title = ""

        for shape in slide.shapes:
            try:
                if getattr(shape, "has_text_frame", False):
                    txt = shape.text_frame.text
                    if txt:
                        slide_text += txt + "\n"
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        img_bytes = shape.image.blob
                        images.append(img_bytes)
                    except Exception as e:
                        print("PPTX画像抽出失敗:", e)
            except Exception as e:
                print("shape処理中エラー:", e)

        slides.append({
            "index": i,
            "title": title,
            "text": slide_text.strip(),
            "images": images
        })
    return slides

# ---------------- ユーティリティ: PDF 抽出 ----------------
def extract_pages_from_pdf_bytes(file_bytes):
    slides = []
    with fitz.open(stream=file_bytes, filetype="pdf") as doc:
        for i, page in enumerate(doc, start=1):
            text = page.get_text().strip()
            images = []
            for img in page.get_images(full=True):
                xref = img[0]
                try:
                    base_image = doc.extract_image(xref)
                    images.append(base_image["image"])
                except Exception as e:
                    print("PDF画像抽出失敗:", e)
            slides.append({
                "index": i,
                "title": "",
                "text": text,
                "images": images
            })
    return slides

# ---------------- フォールバック：簡易クイズ生成 ----------------
def simple_quiz_from_slides(slides, n_questions=5):
    def first_sentence(text):
        s = (text or "").strip().split("\n")[0]
        parts = re.split(r"(?<=。|\.|\?|！|!|？)", s)
        return parts[0].strip() if parts and parts[0].strip() else s[:200].strip()

    s_sorted = sorted(slides, key=lambda s: len(s.get("text","")), reverse=True)
    if not s_sorted:
        return "クイズ生成に失敗しました（素材がありません）"

    out = []
    for i in range(min(n_questions, len(s_sorted))):
        s = s_sorted[i]
        base = first_sentence(s.get("text","") or s.get("title","") or "次の内容について問います。")
        qtext = f"{base} について、最も適切なのはどれか。"
        keywords = re.findall(r"[一-龥ぁ-んァ-ンA-Za-z0-9]+", base)
        kw = keywords[0] if keywords else "（キーワード）"
        correct = f"{kw} の定義に合致する説明"
        wrongs = [f"{kw} に似ているが誤り1", f"{kw} に似ているが誤り2", f"{kw} に似ているが誤り3"]
        opts = [correct] + wrongs
        out.append(f"===問題{i+1}===\n【問題文】\n{qtext}\n【選択肢】\nA. {opts[0]}\nB. {opts[1]}\nC. {opts[2]}\nD. {opts[3]}\n【正解】\nA\n【解説】\n参照スライド：スライド {s.get('index')}（タイトル：'{s.get('title')}'), 該当箇所：'{first_sentence(s.get('text') or s.get('title') or '')}'\n")
    return "\n".join(out)

# ---------------- モデル呼び出し（堅牢化） ----------------
def generate_quiz_from_slides(slides):
    """
    モデル呼び出しの堅牢化 + フォールバック。
    """
    # 必要なら model を初期化
    if not ensure_model():
        print("モデル未利用フォールバック：モデル利用不可のため簡易クイズを使用します。")
        return simple_quiz_from_slides(slides)

    prompt_header = (
        "以下はスライド単位の材料です。各スライドは番号とタイトルを持ちます。"
        "以下の資料内容に基づき、大学生向けの四択クイズを5問作ってください。\n"
        "【作成方針】\n"
        "1. 資料の中から重要な用語・概念・理論・出来事などを選び、その内容を理解しているか確認できる問題を作成してください。\n"
        "2. 各問題は「定義・特徴・関連性・因果関係・応用」など、理解を測る内容にしてください。\n"
        "3. 各問題の選択肢は **必ず4つのみ**（A, B, C, D）を作成してください。\n"
        "4. 各選択肢は自然な文章で構成し、1つだけが明確に正解となるようにしてください。\n"
        "5. 解説では「なぜその選択肢が正解なのか」「他の選択肢がなぜ誤りなのか」を丁寧に説明してください。\n"
        "6. 解説をする際は選択肢ごとにそれぞれ作成してください。"
        "7. 各解説ではスライドの何枚目に関連情報があるかも必ず示してください。\n"
        "8. 「資料の右下に書いています」など曖昧な表現は禁止です。\n"
        "9. **必ず**どのスライドの何枚目（スライド番号）に関連情報があるかを明記すること。表記は次の形式に従う：  「参照スライド：スライド n（タイトル：'スライドのタイトル'）、該当箇所：'（スライド上の見出しや段落を短く引用）'」  （例）参照スライド：スライド 4（タイトル：'〜〜'）、該当箇所：'〜〜の定義'。\n"
        "ただし関連情報の表記はまとめて書くとわかりずらいので選択肢ごとに書いてください。（例）「解説正解はAです。～(解説)～　参照スライド：スライド 4（タイトル：'〜〜'）、該当箇所：'〜〜の定義'\n"
        
        "\n"
        "【出力フォーマット】\n"
        "===問題1===\n"
        "【問題文】\n...\n"
        "【選択肢】\nA. ...\nB. ...\nC. ...\nD. ...\n"
        "【正解】\nA\n"
        "【解説】\n...\n\n"
        "===問題2===\n"
        "...\n"
        "===問題3===\n"
        "...\n"
        "===問題4===\n"
        "...\n"
        "===問題5===\n"
        "...\n"
    )

    parts = [prompt_header]
    for s in slides:
        txt = s["text"][:4000]
        title_line = f"Title: {s['title']}" if s['title'] else "Title: (なし)"
        parts.append(f"===SLIDE {s['index']}===\n{title_line}\n{txt}\n")
        for j, img_bytes in enumerate(s.get("images", []), start=1):
            try:
                img = Image.open(io.BytesIO(img_bytes))
                buf = io.BytesIO()
                img.save(buf, format="PNG")
                img_b64 = base64.b64encode(buf.getvalue()).decode("utf-8")
                parts.append(f"[Image: slide={s['index']} image_index={j}]")
                parts.append({"inline_data": {"mime_type": "image/png", "data": img_b64}})
            except Exception as e:
                print("画像処理失敗:", e)

    parts.append("※解説では必ず「参照スライド：スライド n（タイトル：'...')、該当箇所：'...'」の形式を使ってください。")

    try:
        response = model.generate_content(parts)
    except Exception as e:
        print("=== モデル呼び出し例外 ===")
        print(repr(e))
        print("フォールバックで簡易クイズを生成します。")
        return simple_quiz_from_slides(slides)

    try:
        # 応答抽出の複数パス
        if getattr(response, "candidates", None):
            c0 = response.candidates[0]
            if getattr(c0, "content", None) and getattr(c0.content, "parts", None):
                p0 = c0.content.parts[0]
                if getattr(p0, "text", None):
                    return p0.text
            if getattr(c0, "output_text", None):
                return c0.output_text

        if getattr(response, "text", None):
            return response.text

        print("=== モデル応答（生） ===")
        print(repr(response))
        print("応答から生成テキストを抽出できませんでした。フォールバックを使用します。")
        return simple_quiz_from_slides(slides)

    except Exception as e:
        print("=== 応答処理で例外 ===")
        print(repr(e))
        return simple_quiz_from_slides(slides)

# ---------------- 応答パース ----------------
def parse_quiz_text(response_text):
    questions = []

    # 問題ブロック毎に分割（===問題1=== 等）
    blocks = re.split(r"===\s*問題\s*\d+\s*===", response_text)
    for block in blocks:
        block = block.strip()
        if not block:
            continue

        # 問題文
        m_q = re.search(r"【問題文】\s*(.*?)\s*【選択肢】", block, re.S)
        if not m_q:
            m_q = re.search(r"(?:問題文|問題)\s*[:：]?\s*(.*?)\s*(?:選択肢|A\.)", block, re.S)
        if not m_q:
            continue
        question_text = m_q.group(1).strip()

        # 選択肢
        m_opts = re.search(r"【選択肢】\s*(.*?)(?:\s*【正解】|\s*正解\s*:|\s*【解説】|\Z)", block, re.S)
        if not m_opts:
            continue
        opts_block = m_opts.group(1).strip()

        opts = {}
        for opt_label in ["A", "B", "C", "D"]:
            pattern = rf"{opt_label}\.\s*(.*?)(?=(?:\n[A-D]\.|$))"
            m = re.search(pattern, opts_block, re.S)
            if m:
                opts[opt_label] = m.group(1).strip()
            else:
                opts[opt_label] = ""

        m_ans = re.search(r"【正解】\s*([A-D])", block)
        if not m_ans:
            m_ans = re.search(r"正解\s*[:：]?\s*([A-D])", block)
        answer = m_ans.group(1).strip() if m_ans else ""

        m_exp = re.search(r"【解説】\s*(.*)", block, re.S)
        explanation = m_exp.group(1).strip() if m_exp else ""

        if question_text and all(opts.values()):
            questions.append({
                "question": question_text,
                "options": [f"A. {opts['A']}", f"B. {opts['B']}", f"C. {opts['C']}", f"D. {opts['D']}"],
                "answer": answer,
                "explanation": explanation
            })

    return questions

# ---------------- 復習解析ユーティリティ ----------------
def extract_referenced_slides_from_explanation(explanation_text):
    nums = re.findall(r"スライド\s*(\d+)", explanation_text)
    return sorted({int(n) for n in nums}) if nums else []

def summarize_weak_points(question_list):
    slide_count = {}
    wrong_questions = []
    for q in question_list:
        if not q.get("correct"):
            wrong_questions.append(q)
            refs = extract_referenced_slides_from_explanation(q.get("explanation", ""))
            for r in refs:
                slide_count[r] = slide_count.get(r, 0) + 1

    sorted_slides = sorted(slide_count.items(), key=lambda x: x[1], reverse=True)
    weak_slides = [s for s, cnt in sorted_slides]
    return {
        "total": len(question_list),
        "wrong": len(wrong_questions),
        "weak_slides": weak_slides,
        "wrong_questions": wrong_questions
    }

def generate_study_advice_with_model(wrong_questions, weak_slides, max_tokens=800):
    if not ensure_model():
        print("モデル未利用: モデルでのアドバイス生成はスキップします。")
        return None

    if not wrong_questions:
        return None

    prompt = "以下は受験者が間違えた問題とその解説です。大学生向けに、弱点に合わせた具体的な学習アドバイス（短期プラン + 具体的な復習アクション）を日本語で作ってください。参照スライド番号があればそこを重点に、どう復習すれば理解が深まるかを示してください。\n\n"
    for i, q in enumerate(wrong_questions, start=1):
        prompt += f"---\n間違い{i}:\n問題: {q.get('question')}\n正解: {q.get('answer')}\n解説: {q.get('explanation')[:1000]}\n\n"
    prompt += "\n出力は箇条書き（短い見出し + 具体アクション）でお願いします。"

    try:
        parts = [prompt]
        response = model.generate_content(parts)
        if response.candidates and response.candidates[0].content.parts:
            return response.candidates[0].content.parts[0].text
        if getattr(response, "text", None):
            return response.text
        return None
    except Exception as e:
        print("モデル呼び出し（学習アドバイス）失敗:", repr(e))
        return None

def heuristic_advice(analysis):
    lines = []
    if analysis["wrong"] == 0:
        lines.append("おめでとうございます！全問正解です。理解は十分そうです。")
        lines.append("次のステップ: 応用問題に挑戦して定着を図りましょう（類題を3問以上解くのがおすすめ）。")
        return "\n".join(lines)

    lines.append(f"間違い数: {analysis['wrong']} / {analysis['total']}")
    if analysis["weak_slides"]:
        lines.append("特に復習すべきスライド（出題で対応ができていなかった箇所）:")
        for s in analysis["weak_slides"]:
            lines.append(f" - スライド {s} をもう一度読み、スライド中の定義・式・図の意味をまとめる。")
    else:
        lines.append("解説に参照スライド情報が十分になかったため、間違えた問題の解説文を丁寧に読むことを優先してください。")

    lines.append("\n具体的な学習アクション（短期プラン）:")
    lines.append("1) 間違えた各問題について「なぜ間違えたか」をノートに1行でまとめる（原因分析）。")
    lines.append("2) 参照スライドをテキスト化して、キーワード5つを抜き出す。")
    lines.append("3) キーワードごとに1問ずつ自分で作問して解く（アウトプットで定着）。")
    lines.append("4) 1週間後に同じ問題を再度解き、正答率が上がっているかを確認する。")

    return "\n".join(lines)

# ---------------- ルーティング ----------------
@app.route("/", methods=["GET", "POST"])
def index():
    if request.method == "POST":
        session.clear()

        f = request.files.get("file")
        if not f:
            return "ファイルが選択されていません", 400

        filename = f.filename.lower()
        file_bytes = f.read()
        try:
            if filename.endswith(".pptx"):
                slides = extract_slides_from_pptx_bytes(file_bytes)
            elif filename.endswith(".pdf"):
                slides = extract_pages_from_pdf_bytes(file_bytes)
            else:
                return "対応していないファイル形式です（.pptx または .pdf）", 400
        except Exception as e:
            print("ファイル抽出時エラー:", e)
            return "ファイルの抽出に失敗しました", 500

        response_text = generate_quiz_from_slides(slides)

        print("=== MODEL OUTPUT START ===")
        print(response_text[:5000])
        print("=== MODEL OUTPUT END ===")

        questions = parse_quiz_text(response_text)

        if questions:
            for q in questions:
                q["user_answer"] = None
                q["correct"] = False
                q["ref_slides"] = extract_referenced_slides_from_explanation(q.get("explanation", ""))
            APP_STATE["question_list"] = questions

            session["current_index"] = 0
            session["score"] = 0
            session["total"] = len(questions)

            return redirect(url_for("question_page"))
        else:
            return (
                "<h3>クイズ生成に失敗しました — モデルの応答が期待フォーマットと一致しませんでした。</h3>"
                "<pre>{}</pre>".format(response_text.replace("<", "&lt;").replace(">", "&gt;")),
                500,
            )

    return render_template("index.html")

@app.route("/question", methods=["GET", "POST"])
def question_page():
    question_list = APP_STATE.get("question_list", [])

    current_index = session.get("current_index", 0)
    score = session.get("score", 0)

    if not question_list:
        return redirect(url_for("index"))

    if request.method == "POST":
        selected = request.form.get("answer")
        if selected:
            selected_label = selected.strip()[0].upper()
        else:
            selected_label = None

        correct_label = question_list[current_index]["answer"].strip()

        question_list[current_index]["user_answer"] = selected_label
        is_correct = (selected_label == correct_label)
        question_list[current_index]["correct"] = is_correct

        if is_correct:
            score += 1
            session["score"] = score

        current_index += 1
        if current_index >= len(question_list):
            ##session["question_list"] = question_list
            session["current_index"] = 0
            return redirect(url_for("result_page"))
        else:
            ##session["question_list"] = question_list
            session["current_index"] = current_index

    question_data = question_list[current_index]
    return render_template(
        "question.html",
        question_data=question_data,
        index=current_index + 1,
        total=len(question_list)
    )

@app.route("/result")
def result_page():
    score = session.get("score", 0)
    total = session.get("total", 0)


    return render_template("result.html", score=score, total=total)

@app.route("/review")
def review_page():
    question_list = APP_STATE.get("question_list", [])

    if not question_list:
        return redirect(url_for("index"))

    analysis = summarize_weak_points(question_list)
    return render_template("review.html", question_list=question_list, analysis=analysis)

@app.route("/advice", methods=["POST"])
def advice():
    question_list = APP_STATE.get("question_list", [])

    if not question_list:
        return jsonify({"ok": False, "message": "問題データが見つかりません"}), 400

    analysis = summarize_weak_points(question_list)
    model_advice = generate_study_advice_with_model(analysis["wrong_questions"], analysis["weak_slides"])
    if model_advice:
        return jsonify({"ok": True, "advice": model_advice})
    else:
        return jsonify({"ok": True, "advice": heuristic_advice(analysis)})

@app.route("/update_score", methods=["POST"])
def update_score():
    data = request.get_json()
    if data.get("correct"):
        session["score"] = session.get("score", 0) + 1
    return jsonify({"score": session.get("score", 0)})

# ---------------- 実行 ----------------
if __name__ == "__main__":
    # 開発中は use_reloader=False を指定して一プロセスにすることもできます
    app.run(debug=False, use_reloader=False)
